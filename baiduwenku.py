import requests
import selenium
from selenium import webdriver
import time
import os
from pptx import Presentation  # 生成ppt需要的包
import re
from lxml import etree


import config as cfg
from selenium_helper import SeleniumHelper

def get_html(url, return_type="text"):  # 获取网页页面，返回文本文件或者二进制文件
    r = requests.get(url)
    r.raise_for_status()
    r.encoding = r.apparent_encoding
    if return_type == "text":
        return r.text
    else:
        return r.content


def get_picture_url(url):  # 得到该页面的每页图片的url
    opt = webdriver.FirefoxOptions()  # 设置驱动的浏览器是火狐
    opt.set_headless()
    driver = webdriver.Firefox(options=opt)  # 启动火狐，无界面模式
    driver.get(url)  # 打开url
    time.sleep(3)  # 休息3秒，等待加载
    # xpath寻找图片所在的tag
    img_tags = driver.find_elements_by_xpath("//div[@class='ppt-image-wrap']/img")
    img_urls = []
    for img_tag in img_tags:  # 提取url
        if img_tag.get_attribute("src"):
            img_urls.append(img_tag.get_attribute("src"))
        else:
            img_urls.append(img_tag.get_attribute("data-src"))
    return img_urls  # 返回


def download_pictures(url):  # 下载图片，并存储
    img_urls = get_picture_url(url)  # 得到图片url
    if not os.path.exists("./pictures"):  # 判断是否有存储路径
        os.makedirs("./pictures")
    path_save = []  # 储存图片的名字和路径，在合成ppt时保持正确的顺序
    for i in range(len(img_urls)):
        img = get_html(img_urls[i], return_type="img")  # 得到图片的二进制文件
        with open("./pictures/num_%s.jpg" % str(i), "wb") as f:  # 写出
            path_save.append("./pictures/num_%s.jpg" % str(i))
            f.write(img)
    return path_save  # 返回正确顺序的图片路径和名字


def download_word(driver,url):  # 如果是word， 则调用这个函数
    # opt = webdriver.FirefoxOptions()
    # opt.set_headless()
    # driver = webdriver.Firefox(options=opt)
    # driver.get(url)
    # time.sleep(3)

    all_text = ""
    if re.search("继续阅读", driver.page_source):
        # 下拉到“继续阅读“的位置
        js = "var q=document.documentElement.scrollTop=4500"
        driver.execute_script(js)

        # 点击继续阅读
        hidden_div = driver.find_element_by_xpath("//div[@id='html-reader-go-more']")
        got_btn = driver.find_element_by_xpath("//div[@class='banner-more-btn']/span")
        actions = webdriver.ActionChains(driver)
        actions.move_to_element(hidden_div)
        actions.click(got_btn)
        actions.perform()
        time.sleep(1)

        # 获取页码
        page_tags = driver.find_elements_by_xpath("//div[contains(@class, 'mod reader-page complex')]")
        for i in range(len(page_tags)):  # 循环每页
            print("Total page is %d, current downloading page is %d." % (len(page_tags), i + 1))
            # 定位每页的位置
            driver.execute_script("arguments[0].scrollIntoView();", page_tags[i])
            time.sleep(0.2)
            # 找到相应的p标签，文字 就储存在p标签中
            p_tags = driver.find_elements_by_xpath("//div[contains(@class, 'mod reader-page complex') and contains(@class, 'reader-page-%d')][1]/div/div/div/div/div/div[last()]/div/p" % (i + 1))
            tag_ = True  # 因为有太多的换行，这里处理一下，设置一个判断的变量
            for p_tag in p_tags:
                text = p_tag.text
                if text == "" and tag_ is True:  # 前一个不为空， 当前为空，则换行
                    all_text = all_text + "\n\n"
                    tag_ = False
                elif text == "" and tag_ is False:  # 前一个为空，当前也为空，不变
                    all_text = all_text
                else:  # 其它的就直接拼接
                    all_text = all_text + text
                    tag_ = True


    else:
        p_tags = driver.find_elements_by_xpath("//div[@class='reader-txt-layer']/div/p")
        tag_ = True  # 因为有太多的换行，这里处理一下，设置一个判断的变量
        for p_tag in p_tags:
            text = p_tag.text

            if text == "" and tag_ is True:  # 前一个不为空， 当前为空，则换行
                all_text = all_text + "\n\n"
                tag_ = False
            elif text == "" and tag_ is False:  # 前一个为空，当前也为空，不变
                all_text = all_text
            else:  # 其它的就直接拼接
                all_text = all_text + text
                tag_ = True


    return all_text  # 返回拼接好的字符串


def generate_ppt(picture_paths):  # 如果是ppt， 在上面下载好图片后，这个函数嫁给你图片重新合成ppt
    prs = Presentation('test.pptx')  # 这个是一个ppt模板，里面只有一个图片的占位
    for picture_path in picture_paths:
        # 按照第一个模板创建 一张幻灯片
        oneSlide = prs.slides.add_slide(prs.slide_layouts[0])
        # 获取模板可填充的所有位置
        body_shapes = oneSlide.shapes.placeholders
        for index, body_shape in enumerate(body_shapes):
            if index == 0:
                body_shape.insert_picture(picture_path)

    # 对ppt的修改
    prs.save('baiduwenku.pptx')
    # 这个函数具体的不是很懂，在网上copy的


def write_to_file(browser,url, type,file_name):  # 写入函数
    if type == "ppt":
        path_history = download_pictures(url)
        generate_ppt(path_history)
        print("download successfully!")
    elif type == "word":
        text = download_word(browser,url)
        path=cfg.save_path+'/'+file_name
        with open(path, "w", encoding="utf-8") as f:
            f.write(text)
        print("download successfully!")
    else:
        print("Does not support this type of file now.")

'''
def daokebaba(url):
    html = get_html(url)
    res = re.search("https://www.docin.com/DocinViewer-355106913-144.swf", html)
    if res:
        print(res.group())
    else:
        print("no url!")
    swf_url = res.group()
    swf_res = get_html(swf_url, return_type="img")
    with open("./ddaokebaba.swf", "wb") as f:
        f.write(swf_res)
'''

if __name__ == "__main__":  # 主函数
    # url = "https://wenku.baidu.com/view/dc04b7603069a45177232f60ddccda38376be1f6.html?from=search"
    # url='https://wenku.baidu.com/view/b240e09c4bfe04a1b0717fd5360cba1aa8118c2d.html?from=search'
    # return_type = "word"
    # browser=SeleniumHelper().browser
    #
    # write_to_file(browser,url, return_type,'result.doc')
    # 1， 判断是ppt还是word，
    # 2， 两种不同的文件调用不同的函数
    # url = "https://www.docin.com/p-355106913.html"
    # daokebaba(url)
    pass

